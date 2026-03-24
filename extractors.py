"""Estrattori di testo per chunking e display.

Strumenti di estrazione:
- dots.ocr: UNICO strumento per PDF, immagini, documenti → testo
- Whisper turbo: UNICO strumento per audio → testo
- Video: audio track → Whisper + keyframes → dots.ocr
- Notebook: parsing JSON celle
- Excel: openpyxl
- Presentazioni: python-pptx
- Testo: lettura diretta
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# === dots.ocr (PDF, immagini, documenti) ===

def extract_with_dots_ocr(file_path: Path) -> str:
    """Estrae testo da qualsiasi documento/immagine usando dots.ocr.

    Supporta nativamente: PDF, PNG, JPG, TIFF, BMP, WEBP, ecc.
    https://github.com/rednote-hilab/dots.ocr
    """
    try:
        from dots_ocr import DotsOCR
        ocr = DotsOCR()
        result = ocr.recognize(str(file_path))
        if isinstance(result, list):
            return '\n'.join(str(r) for r in result).strip()
        return str(result).strip()
    except ImportError:
        raise ImportError("pip install dots-ocr  (https://github.com/rednote-hilab/dots.ocr)")
    except Exception as e:
        logger.error(f"dots.ocr fallito per {file_path.name}: {e}")
        raise


# === Whisper (Audio) ===

def extract_audio_text(file_path: Path, model_name: str = "turbo") -> str:
    """Trascrive audio usando Whisper turbo (richiede FFmpeg).

    Nessun embedding model gestisce audio raw → Whisper e' necessario.
    """
    try:
        import whisper
    except ImportError:
        raise ImportError("pip install openai-whisper  (richiede FFmpeg)")

    logger.info(f"Whisper '{model_name}': trascrizione {file_path.name}...")
    model = whisper.load_model(model_name)
    result = model.transcribe(
        str(file_path),
        language=None,          # Auto-detect lingua
        task="transcribe",
        verbose=False,
    )
    return result["text"].strip()


# === Video (audio + keyframes) ===

def extract_video_text(file_path: Path, config: dict = None) -> str:
    """Estrae contenuto da video: audio → Whisper + keyframes → dots.ocr"""
    parts = []
    whisper_model = (config or {}).get('whisper_model', 'turbo')

    # 1. Audio track → Whisper
    try:
        import subprocess
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp:
            tmp_path = tmp.name

        subprocess.run(
            ['ffmpeg', '-i', str(file_path), '-vn', '-acodec', 'pcm_s16le',
             '-ar', '16000', '-ac', '1', tmp_path, '-y'],
            capture_output=True, check=True
        )
        audio_text = extract_audio_text(Path(tmp_path), whisper_model)
        if audio_text:
            parts.append(f"[AUDIO TRANSCRIPT]\n{audio_text}")
        Path(tmp_path).unlink(missing_ok=True)
    except Exception as e:
        logger.warning(f"Estrazione audio da video fallita: {e}")

    # 2. Keyframes → dots.ocr
    try:
        import cv2
        import tempfile
        cap = cv2.VideoCapture(str(file_path))
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = cap.get(cv2.CAP_PROP_FPS) or 30

        # 1 frame ogni 30 secondi, max 10 frame
        interval = int(fps * 30)
        frame_indices = list(range(0, total_frames, max(interval, 1)))[:10]

        for fi in frame_indices:
            cap.set(cv2.CAP_PROP_POS_FRAMES, fi)
            ret, frame = cap.read()
            if not ret:
                continue
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                cv2.imwrite(tmp.name, frame)
                try:
                    frame_text = extract_with_dots_ocr(Path(tmp.name))
                    if frame_text:
                        timestamp = fi / fps
                        parts.append(f"[FRAME {timestamp:.0f}s]\n{frame_text}")
                except Exception:
                    pass
                Path(tmp.name).unlink(missing_ok=True)
        cap.release()
    except ImportError:
        logger.warning("opencv-python-headless non installato per estrazione frame video")
    except Exception as e:
        logger.warning(f"Estrazione frame video fallita: {e}")

    return '\n\n---\n\n'.join(parts).strip() if parts else f"[Video: {file_path.name}]"


# === Notebook ===

def extract_notebook_text(file_path: Path) -> str:
    """Estrae testo strutturato da Jupyter Notebook (codice, markdown, output)."""
    import json
    nb = json.loads(file_path.read_text(encoding='utf-8'))
    parts = []

    for idx, cell in enumerate(nb.get('cells', []), 1):
        cell_type = cell.get('cell_type', 'raw')
        source = ''.join(cell.get('source', []))

        if not source.strip() and cell_type != 'code':
            continue

        if cell_type == 'markdown':
            parts.append(f"[CELL {idx} | MARKDOWN]\n{source}")
        elif cell_type == 'code':
            parts.append(f"[CELL {idx} | CODE]\n```python\n{source}\n```")
            for out in cell.get('outputs', []):
                otype = out.get('output_type', '')
                if otype == 'stream':
                    text = ''.join(out.get('text', []))
                    if text.strip():
                        parts.append(f"[CELL {idx} | OUTPUT]\n{text.strip()}")
                elif otype in ('execute_result', 'display_data'):
                    text = ''.join(out.get('data', {}).get('text/plain', []))
                    if text.strip():
                        parts.append(f"[CELL {idx} | RESULT]\n{text.strip()}")
                elif otype == 'error':
                    err = f"{out.get('ename', 'Error')}: {out.get('evalue', '')}"
                    parts.append(f"[CELL {idx} | ERROR]\n{err}")
        else:
            parts.append(f"[CELL {idx} | {cell_type.upper()}]\n{source}")

    return '\n\n---\n\n'.join(parts).strip()


# === Excel ===

def extract_excel_text(file_path: Path) -> str:
    """Estrae testo da Excel (.xlsx)."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    texts = []
    for sheet in wb:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = ' | '.join(str(c) for c in row if c is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            texts.append(f"[SHEET: {sheet.title}]\n" + '\n'.join(rows))
    return '\n\n---\n\n'.join(texts).strip()


# === Presentazioni ===

def extract_presentation_text(file_path: Path) -> str:
    """Estrae testo da PowerPoint (.pptx) - slide + note."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    parts = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        if slide_texts or notes_text:
            part = f"[SLIDE {idx}]\n" + '\n'.join(slide_texts)
            if notes_text:
                part += f"\n[NOTE]\n{notes_text}"
            parts.append(part)

    return '\n\n---\n\n'.join(parts).strip()


# === Dispatcher ===

def extract_text(file_path: Path, config: dict = None) -> tuple[str, str | None]:
    """Dispatcher: estrae testo da qualsiasi formato supportato.

    Returns:
        (testo, errore) - se errore e' None, l'estrazione e' riuscita
    """
    ext = file_path.suffix.lower()
    cfg = config or {}

    audio_ext = set(cfg.get('extensions', {}).get('audio', ['.mp3', '.m4a', '.wav', '.ogg', '.flac']))
    video_ext = set(cfg.get('extensions', {}).get('video', ['.mp4', '.mkv', '.avi']))
    notebook_ext = set(cfg.get('extensions', {}).get('notebook', ['.ipynb']))
    excel_ext = set(cfg.get('extensions', {}).get('excel', ['.xlsx']))
    presentation_ext = set(cfg.get('extensions', {}).get('presentation', ['.pptx']))
    pdf_ext = set(cfg.get('extensions', {}).get('pdf', ['.pdf']))
    image_ext = set(cfg.get('extensions', {}).get('image', ['.png', '.jpg', '.jpeg']))

    try:
        # dots.ocr per PDF e immagini
        if ext in pdf_ext or ext in image_ext:
            return extract_with_dots_ocr(file_path), None

        # Whisper per audio
        if ext in audio_ext:
            return extract_audio_text(file_path, cfg.get('whisper_model', 'turbo')), None

        # Video: audio + keyframes
        if ext in video_ext:
            return extract_video_text(file_path, cfg), None

        # Formati strutturati
        if ext in notebook_ext:
            return extract_notebook_text(file_path), None
        if ext in excel_ext:
            return extract_excel_text(file_path), None
        if ext in presentation_ext:
            return extract_presentation_text(file_path), None

        # Testo semplice
        return file_path.read_text(encoding='utf-8', errors='ignore').strip(), None

    except ImportError as e:
        return "", f"Dipendenza mancante: {e}"
    except Exception as e:
        logger.error(f"Errore estrazione {file_path}: {e}")
        return "", str(e)
